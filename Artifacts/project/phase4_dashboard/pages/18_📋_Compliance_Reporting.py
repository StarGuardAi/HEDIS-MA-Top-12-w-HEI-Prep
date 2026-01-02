"""
Automated Compliance Reporting System

Handles:
- CMS Star Rating Submission
- NCQA HEDIS Reporting
- Internal Compliance Reports
- Audit Support Documentation
"""

import streamlit as st
import pandas as pd
from datetime import datetime, timedelta
from utils.compliance_reporting import (
    ComplianceReportingEngine,
    ReportType,
    ReportStatus
)
from utils.database import get_db_connection
from src.ui.layout import render_page_footer, render_sidebar_footer, render_header, render_starguard_header
import io
import base64

# Check for optional export libraries
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

st.set_page_config(
    page_title="Compliance Reporting",
    page_icon="📋",
    layout="wide"
)

# Rename "app" to "Home" in sidebar navigation
st.markdown("""
<style>
/* Method 1: Replace text content with CSS */
[data-testid="stSidebarNav"] ul li:first-child a {
    font-size: 0 !important;  /* Hide original "app" text */
    background: linear-gradient(135deg, rgba(139,122,184,0.3), rgba(111,95,150,0.3)) !important;
    padding: 0.75rem 1rem !important;
    border-radius: 8px !important;
    border: 1px solid rgba(255,255,255,0.2) !important;
    margin-bottom: 0.5rem !important;
}

[data-testid="stSidebarNav"] ul li:first-child a::before {
    content: "🏠 Home";  /* Replace with "Home" */
    font-size: 1.1rem !important;
    color: white !important;
    font-weight: 700 !important;
}

[data-testid="stSidebarNav"] ul li:first-child a:hover {
    background: linear-gradient(135deg, rgba(139,122,184,0.5), rgba(111,95,150,0.5)) !important;
}

/* Method 2: Also target by href */
[data-testid="stSidebarNav"] a[href="/"],
[data-testid="stSidebarNav"] a[href="./"] {
    font-size: 0 !important;
    background: linear-gradient(135deg, rgba(139,122,184,0.3), rgba(111,95,150,0.3)) !important;
    padding: 0.75rem 1rem !important;
    border-radius: 8px !important;
}

[data-testid="stSidebarNav"] a[href="/"]::before,
[data-testid="stSidebarNav"] a[href="./"]::before {
    content: "🏠 Home";
    font-size: 1.1rem !important;
    color: white !important;
    font-weight: 700 !important;
}

/* ========== PURPLE SIDEBAR THEME (Match Home Page) ========== */
/* Match the StarGuard AI header purple gradient */
[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #4A3D6F 0%, #6F5F96 100%) !important;
}

[data-testid="stSidebar"] > div:first-child {
    background: linear-gradient(180deg, #4A3D6F 0%, #6F5F96 100%) !important;
}

/* ========== ALL SIDEBAR TEXT WHITE ========== */
/* Force ALL text in sidebar to be white */
[data-testid="stSidebar"] * {
    color: #FFFFFF !important;
}

[data-testid="stSidebar"] p,
[data-testid="stSidebar"] span,
[data-testid="stSidebar"] div,
[data-testid="stSidebar"] a,
[data-testid="stSidebar"] label,
[data-testid="stSidebar"] button {
    color: #FFFFFF !important;
}

/* ========== WHITE "HOME" LABEL ========== */
[data-testid="stSidebarNav"] ul li:first-child a {
    font-size: 0 !important;
    background: rgba(255, 255, 255, 0.2) !important;
    padding: 0.75rem 1rem !important;
    border-radius: 8px !important;
    border: 2px solid rgba(255, 255, 255, 0.3) !important;
    margin-bottom: 0.5rem !important;
}

[data-testid="stSidebarNav"] ul li:first-child a::before {
    content: "🏠 Home" !important;
    font-size: 1.1rem !important;
    color: #FFFFFF !important;
    font-weight: 700 !important;
    display: block !important;
    -webkit-text-fill-color: #FFFFFF !important;
    text-shadow: 0 1px 2px rgba(0, 0, 0, 0.3) !important;
}

/* All sidebar navigation links white */
[data-testid="stSidebarNav"] a {
    color: #FFFFFF !important;
}

[data-testid="stSidebarNav"] a span,
[data-testid="stSidebarNav"] a div,
[data-testid="stSidebarNav"] a p {
    color: #FFFFFF !important;
}

/* "Mobile Optimized" badge - white text */
[data-testid="stSidebar"] .element-container div[data-testid="stMarkdownContainer"] p {
    color: #FFFFFF !important;
}

/* Success/Info boxes in sidebar - white text */
[data-testid="stSidebar"] [data-testid="stSuccess"],
[data-testid="stSidebar"] [data-testid="stInfo"] {
    color: #FFFFFF !important;
    background: rgba(255, 255, 255, 0.15) !important;
    border-color: rgba(255, 255, 255, 0.3) !important;
}

[data-testid="stSidebar"] [data-testid="stSuccess"] *,
[data-testid="stSidebar"] [data-testid="stInfo"] * {
    color: #FFFFFF !important;
}

/* View less/more links - white */
[data-testid="stSidebar"] button {
    color: #FFFFFF !important;
}

/* Mobile responsive */
@media (max-width: 768px) {
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #4A3D6F 0%, #6F5F96 100%) !important;
    }
    
    [data-testid="stSidebar"] * {
        color: #FFFFFF !important;
    }
}

/* ========== HEADER CONTAINER STYLES (Match Home Page) ========== */
.header-container {
    background: linear-gradient(180deg, #4A3D6F 0%, #6F5F96 100%);
    padding: 0.5rem 1rem;
    border-radius: 10px;
    margin-top: 1.5rem;
    margin-bottom: 0.3rem;
    text-align: center;
    display: block !important;
    visibility: visible !important;
    opacity: 1 !important;
    box-shadow: 0 4px 12px rgba(74, 61, 111, 0.25);
    max-width: 100%;
}

.header-title {
    color: white !important;
    font-weight: 700;
    font-size: 1.25rem;
    margin-bottom: 0.4rem;
    display: block !important;
    line-height: 1.5;
    letter-spacing: 0.3px;
}

.header-subtitle {
    color: #E8D4FF !important;
    font-size: 0.9rem;
    font-style: italic;
    display: block !important;
    line-height: 1.4;
    opacity: 0.95;
}
</style>
""", unsafe_allow_html=True)


# Responsive Header - Adapts to Desktop/Mobile (Match Home Page)
st.markdown("""
<div class="header-container">
    <div class="header-title">⭐ StarGuard AI | Turning Data Into Stars</div>
    <div class="header-subtitle">Powered by Predictive Analytics & Machine Learning</div>
</div>
""", unsafe_allow_html=True)

# Spacing fix
st.markdown('<style>div.block-container{padding-top:1rem!important}h1{margin-top:0.5rem!important}</style>', unsafe_allow_html=True)

# Improved compact CSS - READABLE fonts, reduced spacing only
st.markdown("""
<style>
.main .block-container { 
    padding-top: 1rem !important; 
    padding-bottom: 1rem !important; 
    padding-left: 1rem !important; 
    padding-right: 1rem !important; 
    max-width: 100% !important; 
}

/* Section spacing - REDUCE GAPS between sections */
h1 { 
    font-size: 1.8rem !important; 
    margin-top: 0.8rem !important; 
    margin-bottom: 0.5rem !important; 
    line-height: 1.2 !important; 
}

h2 { 
    font-size: 1.4rem !important; 
    margin-top: 0.6rem !important; 
    margin-bottom: 0.4rem !important; 
    line-height: 1.2 !important; 
}

h3 { 
    font-size: 1.1rem !important; 
    margin-top: 0.5rem !important; 
    margin-bottom: 0.3rem !important; 
    line-height: 1.2 !important; 
}

/* Reduce spacing between elements */
.element-container { margin-bottom: 0.4rem !important; }
.stMarkdown { margin-bottom: 0.4rem !important; }

/* Readable metric fonts */
[data-testid="stMetricValue"] { font-size: 1.6rem !important; }
[data-testid="stMetricLabel"] { font-size: 0.95rem !important; padding-bottom: 0.3rem !important; }
[data-testid="metric-container"] { padding: 0.7rem !important; }

/* Chart and data spacing */
.stPlotlyChart { margin-bottom: 0.6rem !important; }
.stDataFrame { margin-bottom: 0.6rem !important; }

/* Column spacing */
[data-testid="column"] { padding: 0.3rem !important; }

/* Interactive elements */
[data-testid="stExpander"] { margin-bottom: 0.5rem !important; }
[data-testid="stTabs"] { margin-bottom: 0.6rem !important; }
.stTabs [data-baseweb="tab-list"] { gap: 0.3rem !important; }
.stTabs [data-baseweb="tab"] { 
    padding: 0.5rem 1rem !important; 
    font-size: 0.95rem !important; 
}

/* Buttons - keep readable */
.stButton > button { 
    padding: 0.6rem 1.2rem !important; 
    font-size: 0.95rem !important; 
}

/* Form inputs */
.stSelectbox, .stTextInput, .stNumberInput { margin-bottom: 0.4rem !important; }

/* Alerts - keep readable */
.stAlert { 
    padding: 0.7rem !important; 
    margin-bottom: 0.5rem !important; 
    font-size: 0.95rem !important; 
}

/* Reduce gaps between blocks */
div[data-testid="stVerticalBlock"] > div { gap: 0.4rem !important; }

/* Horizontal rules */
hr { margin: 0.6rem 0 !important; }

/* Mobile adjustments - Match Home page formatting */
@media (max-width: 768px) {
    .header-container {
        padding: 0.6rem 0.8rem;
        border-radius: 6px;
        margin-top: 1rem;
        margin-bottom: 0.5rem;
        box-shadow: 0 2px 4px rgba(74, 61, 111, 0.15);
    }
    
    .header-title {
        font-size: 0.9rem;
        margin-bottom: 0.25rem;
        line-height: 1.3;
        font-weight: 600;
    }
    
    .header-subtitle {
        font-size: 0.65rem;
        line-height: 1.2;
    }
    
    /* Mobile spacing - tighter */
    div.block-container {
        padding-top: 2rem !important;
        padding-left: 1rem !important;
        padding-right: 1rem !important;
    }
    
    h1 {
        margin-top: 0.5rem !important;
        font-size: 1.5rem !important;
        line-height: 1.3;
        text-align: center !important;
    }
    
    h2 {
        margin-top: 0.75rem !important;
        font-size: 1.25rem !important;
        line-height: 1.3;
        text-align: center !important;
    }
    
    h3 {
        font-size: 1.1rem !important;
        line-height: 1.3;
        text-align: center !important;
    }
    
    h4, h5, h6 {
        text-align: center !important;
    }
    
    /* Center align markdown headers on mobile */
    div[data-testid="stMarkdownContainer"] h1,
    div[data-testid="stMarkdownContainer"] h2,
    div[data-testid="stMarkdownContainer"] h3,
    div[data-testid="stMarkdownContainer"] h4,
    div[data-testid="stMarkdownContainer"] h5,
    div[data-testid="stMarkdownContainer"] h6 {
        text-align: center !important;
    }
    
    /* Center align metrics on mobile */
    [data-testid="stMetric"],
    [data-testid="stMetricValue"],
    [data-testid="stMetricLabel"],
    [data-testid="stMetricDelta"],
    [data-testid="metric-container"],
    .compact-metric-card,
    .kpi-card {
        text-align: center !important;
    }
    
    /* Mobile columns - stack vertically */
    [data-testid="column"] {
        width: 100% !important;
        flex: 1 1 100% !important;
        padding: 0.2rem !important;
    }
    
    /* Mobile buttons - full width */
    button[kind="primary"],
    button[kind="secondary"] {
        width: 100% !important;
        margin-bottom: 0.5rem !important;
    }
    
    /* Mobile metrics - smaller */
    [data-testid="stMetricValue"] {
        font-size: 1.5rem !important;
    }
    
    [data-testid="stMetricLabel"] {
        font-size: 0.85rem !important;
    }
    
    /* Mobile tables - horizontal scroll */
    .stDataFrame {
        overflow-x: auto !important;
    }
    
    /* Mobile tabs - stack vertically to eliminate horizontal scrolling */
    [data-testid="stTabs"] {
        overflow-x: visible !important;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        flex-direction: column !important;
        width: 100% !important;
        gap: 0.5rem !important;
        overflow-x: visible !important;
    }
    
    .stTabs [data-baseweb="tab"] {
        width: 100% !important;
        flex: 1 1 100% !important;
    }
    
    /* Wrap Plotly chart titles on mobile */
    .js-plotly-plot .gtitle,
    .plotly .gtitle,
    .js-plotly-plot .xtitle,
    .plotly .xtitle {
        word-wrap: break-word !important;
        white-space: normal !important;
        max-width: 100% !important;
        overflow-wrap: break-word !important;
        hyphens: auto !important;
    }
    
    /* Ensure chart titles wrap - target SVG text elements */
    .js-plotly-plot .gtitle text,
    .plotly .gtitle text {
        word-wrap: break-word !important;
        white-space: normal !important;
    }
}
</style>
""", unsafe_allow_html=True)

# Page header (already rendered above)

st.title("📋 Automated Compliance Reporting")
st.markdown("Generate, manage, and submit compliance reports with full audit trail")

# Initialize reporting engine
@st.cache_resource
def get_reporting_engine():
    db = get_db_connection()
    return ComplianceReportingEngine(db_connection=db)

reporting = get_reporting_engine()

# Main tabs
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "📊 Generate Reports",
    "📁 Report Library",
    "✅ Approvals",
    "📤 Submissions",
    "🔍 Audit Trail"
])

with tab1:
    st.header("Generate New Report")
    
    report_type = st.selectbox(
        "Select Report Type",
        [
            "CMS Star Rating Submission",
            "NCQA HEDIS Reporting",
            "Internal QA Report",
            "Audit Support Documentation",
            "Board Report"
        ]
    )
    
    generated_by = st.text_input("Generated By", value="Current User")
    
    if report_type == "CMS Star Rating Submission":
        st.subheader("CMS Star Rating Submission")
        
        col1, col2 = st.columns(2)
        with col1:
            data_as_of = st.date_input("Data As Of Date", value=datetime.now().date(), format="MM/DD/YYYY")
        with col2:
            assumptions = st.text_area("Assumptions/Notes", height=100)
        
        if st.button("Generate CMS Star Rating Report", type="primary"):
            with st.spinner("Generating report..."):
                report_data = reporting.create_cms_star_rating_report(
                    data_as_of=datetime.combine(data_as_of, datetime.min.time()),
                    generated_by=generated_by,
                    assumptions={'notes': assumptions} if assumptions else None
                )
                
                st.success(f"Report generated: {report_data['metadata'].report_id}")
                st.session_state['last_report_id'] = report_data['metadata'].report_id
                st.session_state['last_report_data'] = report_data
    
    elif report_type == "NCQA HEDIS Reporting":
        st.subheader("NCQA HEDIS Reporting")
        
        col1, col2 = st.columns(2)
        with col1:
            reporting_period = st.selectbox(
                "Reporting Period",
                ["MY2024", "MY2025", "MY2026"]
            )
        with col2:
            hybrid_methodology = st.checkbox("Use Hybrid Methodology", value=True)
        
        assumptions = st.text_area("Assumptions/Notes", height=100)
        
        if st.button("Generate NCQA HEDIS Report", type="primary"):
            with st.spinner("Generating report..."):
                report_data = reporting.create_ncqa_hedis_report(
                    reporting_period=reporting_period,
                    generated_by=generated_by,
                    hybrid_methodology=hybrid_methodology,
                    assumptions={'notes': assumptions} if assumptions else None
                )
                
                st.success(f"Report generated: {report_data['metadata'].report_id}")
                st.session_state['last_report_id'] = report_data['metadata'].report_id
                st.session_state['last_report_data'] = report_data
    
    elif report_type == "Internal QA Report":
        st.subheader("Internal QA Report")
        
        col1, col2 = st.columns(2)
        with col1:
            report_month = st.selectbox(
                "Report Month",
                [f"{datetime.now().year}-{i:02d}" for i in range(1, 13)]
            )
        with col2:
            include_variance = st.checkbox("Include Variance Explanations", value=True)
        
        assumptions = st.text_area("Assumptions/Notes", height=100)
        
        if st.button("Generate Internal QA Report", type="primary"):
            with st.spinner("Generating report..."):
                report_data = reporting.create_internal_qa_report(
                    report_month=report_month,
                    generated_by=generated_by,
                    include_variance=include_variance,
                    assumptions={'notes': assumptions} if assumptions else None
                )
                
                st.success(f"Report generated: {report_data['metadata'].report_id}")
                st.session_state['last_report_id'] = report_data['metadata'].report_id
                st.session_state['last_report_data'] = report_data
    
    elif report_type == "Audit Support Documentation":
        st.subheader("Audit Support Documentation")
        
        col1, col2 = st.columns(2)
        with col1:
            member_id = st.text_input("Member ID (optional)", value="")
        with col2:
            measure = st.selectbox(
                "Measure",
                ["HbA1c Testing", "Blood Pressure Control", "Breast Cancer Screening"]
            )
        
        assumptions = st.text_area("Assumptions/Notes", height=100)
        
        if st.button("Generate Audit Support Report", type="primary"):
            with st.spinner("Generating report..."):
                report_data = reporting.create_audit_support_report(
                    member_id=member_id if member_id else None,
                    measure=measure,
                    generated_by=generated_by,
                    assumptions={'notes': assumptions} if assumptions else None
                )
                
                st.success(f"Report generated: {report_data['metadata'].report_id}")
                st.session_state['last_report_id'] = report_data['metadata'].report_id
                st.session_state['last_report_data'] = report_data
    
    elif report_type == "Board Report":
        st.subheader("Board Report")
        
        quarter = st.selectbox(
            "Quarter",
            ["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024", "Q1 2025"]
        )
        
        assumptions = st.text_area("Assumptions/Notes", height=100)
        
        if st.button("Generate Board Report", type="primary"):
            with st.spinner("Generating report..."):
                report_data = reporting.create_board_report(
                    quarter=quarter,
                    generated_by=generated_by,
                    assumptions={'notes': assumptions} if assumptions else None
                )
                
                st.success(f"Report generated: {report_data['metadata'].report_id}")
                st.session_state['last_report_id'] = report_data['metadata'].report_id
                st.session_state['last_report_data'] = report_data
    
    # Display last generated report
    if 'last_report_data' in st.session_state:
        st.divider()
        st.subheader("📄 Generated Report Preview")
        
        report_data = st.session_state['last_report_data']
        metadata = report_data['metadata']
        
        # Report summary
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Report ID", metadata.report_id[:20] + "...")
        with col2:
            st.metric("Status", metadata.status.value.title())
        with col3:
            st.metric("Version", metadata.version)
        with col4:
            st.metric("Generated", metadata.generated_at.strftime("%m/%d/%Y"))
        
        # Report-specific preview
        if metadata.report_type == ReportType.CMS_STAR_RATING:
            st.markdown("### Measure Data")
            st.dataframe(report_data['measure_data'], use_container_width=True)
            
            st.markdown("### Domain Scores")
            domain_df = pd.DataFrame([
                {'Domain': k, 'Score': f"{v:.2f}"}
                for k, v in report_data['domain_scores'].items()
            ])
            st.dataframe(domain_df, use_container_width=True)
            
            st.metric("Overall Star Rating", f"{report_data['overall_rating']:.2f}")
            
            st.markdown("### Validation Results")
            validation_df = pd.DataFrame(report_data['validation_results'])
            st.dataframe(validation_df, use_container_width=True)
            
            if report_data['submission_ready']:
                st.success("✅ Report is ready for submission")
            else:
                st.warning("⚠️ Report requires review before submission")
        
        elif metadata.report_type == ReportType.NCQA_HEDIS:
            st.markdown("### Measure Results")
            for measure_result in report_data['measure_results']:
                with st.expander(f"{measure_result['measure']} - {measure_result['rate']:.2f}%"):
                    st.write(f"**Numerator**: {measure_result['numerator']}")
                    st.write(f"**Denominator**: {measure_result['denominator']}")
                    st.write(f"**Rate**: {measure_result['rate']:.2f}%")
                    st.write(f"**Hybrid Methodology**: {measure_result['hybrid_methodology']}")
                    if measure_result['hybrid_methodology']:
                        st.write(f"**Sample Size**: {measure_result['sample_size']}")
        
        elif metadata.report_type == ReportType.INTERNAL_QA:
            st.markdown("### Measure Performance")
            st.dataframe(report_data['measure_data'], use_container_width=True)
            
            if not report_data['corrective_actions'].empty:
                st.markdown("### Corrective Actions")
                st.dataframe(report_data['corrective_actions'], use_container_width=True)
        
        elif metadata.report_type == ReportType.AUDIT_SUPPORT:
            st.markdown("### Member Documentation")
            member_docs = report_data['member_documentation']
            st.json(member_docs)
            
            st.markdown("### Process Documentation")
            st.json(report_data['process_documentation'])
        
        elif metadata.report_type == ReportType.BOARD_REPORT:
            st.markdown("### Executive Summary")
            summary = report_data['executive_summary']
            st.metric("Overall Star Rating", f"{summary['overall_star_rating']:.1f}")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("#### Key Highlights")
                for highlight in summary['key_highlights']:
                    st.write(f"✅ {highlight}")
            with col2:
                st.markdown("#### Risks")
                for risk in summary['risks']:
                    st.write(f"⚠️ {risk}")
            
            st.markdown("### Measure Summary")
            st.dataframe(report_data['measure_summary'], use_container_width=True)
        
        # Export options
        st.markdown("### Export Options")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            # CSV export
            df = reporting.export_to_dataframe(metadata.report_id)
            if not df.empty:
                csv = df.to_csv(index=False)
                st.download_button(
                    "📊 Download CSV",
                    csv,
                    file_name=f"{metadata.report_id}.csv",
                    mime="text/csv"
                )
        
        with col2:
            # Excel export
            if not df.empty:
                excel_buffer = io.BytesIO()
                with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                    df.to_excel(writer, sheet_name='Data', index=False)
                    # Add formulas sheet for Excel
                    formulas_df = pd.DataFrame({
                        'Formula': ['=SUM(Data!B:B)', '=AVERAGE(Data!B:B)'],
                        'Description': ['Total', 'Average']
                    })
                    formulas_df.to_excel(writer, sheet_name='Formulas', index=False)
                excel_buffer.seek(0)
                st.download_button(
                    "📗 Download Excel",
                    excel_buffer,
                    file_name=f"{metadata.report_id}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
        
        with col3:
            # PDF export
            if PDF_AVAILABLE:
                if not df.empty:
                    try:
                        pdf_buffer = io.BytesIO()
                        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
                        styles = getSampleStyleSheet()
                        story = []
                        
                        # Title
                        title = Paragraph(f"<b>{metadata.report_id}</b>", styles['Title'])
                        story.append(title)
                        story.append(Spacer(1, 12))
                        
                        # Report metadata
                        generated_date_str = metadata.generated_at.strftime("%m/%d/%Y") if hasattr(metadata.generated_at, 'strftime') else str(metadata.generated_at)
                        meta_text = f"<b>Report Type:</b> {metadata.report_type}<br/>"
                        meta_text += f"<b>Generated:</b> {generated_date_str}<br/>"
                        meta_text += f"<b>Status:</b> {metadata.status}"
                        meta_para = Paragraph(meta_text, styles['Normal'])
                        story.append(meta_para)
                        story.append(Spacer(1, 12))
                        
                        # Data table
                        table_data = [df.columns.tolist()] + df.values.tolist()
                        table = Table(table_data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(table)
                        
                        doc.build(story)
                        pdf_buffer.seek(0)
                        
                        st.download_button(
                            "📄 Download PDF",
                            pdf_buffer,
                            file_name=f"{metadata.report_id}.pdf",
                            mime="application/pdf"
                        )
                    except Exception as e:
                        st.error(f"PDF generation failed: {str(e)}")
                else:
                    st.info("📄 No data to export")
            else:
                st.info("📄 PDF export requires reportlab library")
                st.caption("Install: pip install reportlab")
        
        with col4:
            # PowerPoint export
            if PPTX_AVAILABLE:
                if not df.empty:
                    try:
                        pptx_buffer = io.BytesIO()
                        prs = Presentation()
                        
                        # Title slide
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        title.text = metadata.report_id
                        generated_date_str = metadata.generated_at.strftime("%m/%d/%Y") if hasattr(metadata.generated_at, 'strftime') else str(metadata.generated_at)
                        subtitle.text = f"Generated: {generated_date_str}\nStatus: {metadata.status}"
                        
                        # Data slide
                        blank_slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(blank_slide_layout)
                        
                        # Add title
                        left = Inches(0.5)
                        top = Inches(0.5)
                        width = Inches(9)
                        height = Inches(0.5)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = "Report Data"
                        tf.paragraphs[0].font.size = Pt(24)
                        tf.paragraphs[0].font.bold = True
                        
                        # Add table (simplified - show first 10 rows)
                        table_data = df.head(10)
                        rows = len(table_data) + 1
                        cols = len(table_data.columns)
                        
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(5)
                        
                        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                        table = table_shape.table
                        
                        # Header row
                        for col_idx, col_name in enumerate(table_data.columns):
                            cell = table.cell(0, col_idx)
                            cell.text = str(col_name)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                            cell.text_frame.paragraphs[0].font.bold = True
                        
                        # Data rows
                        for row_idx, (_, row) in enumerate(table_data.iterrows(), start=1):
                            for col_idx, value in enumerate(row):
                                cell = table.cell(row_idx, col_idx)
                                cell.text = str(value)
                        
                        prs.save(pptx_buffer)
                        pptx_buffer.seek(0)
                        
                        st.download_button(
                            "📊 Download PowerPoint",
                            pptx_buffer,
                            file_name=f"{metadata.report_id}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"PowerPoint generation failed: {str(e)}")
                else:
                    st.info("📊 No data to export")
            else:
                st.info("📊 PowerPoint export requires python-pptx library")
                st.caption("Install: pip install python-pptx")

with tab2:
    st.header("Report Library")
    
    # Filters
    col1, col2 = st.columns(2)
    with col1:
        filter_type = st.selectbox(
            "Filter by Type",
            ["All"] + [rt.value.replace("_", " ").title() for rt in ReportType]
        )
    with col2:
        filter_status = st.selectbox(
            "Filter by Status",
            ["All"] + [rs.value.title() for rs in ReportStatus]
        )
    
    # Get reports
    report_type_filter = None
    if filter_type != "All":
        for rt in ReportType:
            if rt.value.replace("_", " ").title() == filter_type:
                report_type_filter = rt
                break
    
    status_filter = None
    if filter_status != "All":
        for rs in ReportStatus:
            if rs.value.title() == filter_status:
                status_filter = rs
                break
    
    reports = reporting.list_reports(
        report_type=report_type_filter,
        status=status_filter
    )
    
    st.metric("Total Reports", len(reports))
    
    if reports:
        # Display reports in a table
        reports_df = pd.DataFrame(reports)
        st.dataframe(
            reports_df,
            use_container_width=True,
            hide_index=True
        )
        
        # Select report to view
        selected_report_id = st.selectbox(
            "Select Report to View",
            [r['report_id'] for r in reports]
        )
        
        if selected_report_id:
            report_data = reporting.reports.get(selected_report_id)
            if report_data:
                st.subheader("Report Details")
                metadata = report_data['metadata']
                
                # Display metadata
                col1, col2 = st.columns(2)
                with col1:
                    st.write(f"**Report ID**: {metadata.report_id}")
                    st.write(f"**Report Name**: {metadata.report_name}")
                    st.write(f"**Type**: {metadata.report_type.value}")
                    st.write(f"**Status**: {metadata.status.value}")
                with col2:
                    st.write(f"**Generated By**: {metadata.generated_by}")
                    generated_at_str = metadata.generated_at.strftime("%m/%d/%Y") if hasattr(metadata.generated_at, 'strftime') else str(metadata.generated_at)
                    data_as_of_str = metadata.data_as_of.strftime("%m/%d/%Y") if hasattr(metadata.data_as_of, 'strftime') else str(metadata.data_as_of)
                    st.write(f"**Generated At**: {generated_at_str}")
                    st.write(f"**Data As Of**: {data_as_of_str}")
                    st.write(f"**Version**: {metadata.version}")
                
                # Calculations log
                st.markdown("### Calculations Log")
                for calc in metadata.calculations_log:
                    st.write(f"- {calc}")
                
                # Export
                df = reporting.export_to_dataframe(selected_report_id)
                if not df.empty:
                    csv = df.to_csv(index=False)
                    st.download_button(
                        "Download Report",
                        csv,
                        file_name=f"{metadata.report_id}.csv",
                        mime="text/csv"
                    )
    else:
        st.info("No reports found matching criteria")

with tab3:
    st.header("Approval Workflow")
    
    # Get pending approval reports
    pending_reports = reporting.list_reports(status=ReportStatus.PENDING_APPROVAL)
    
    if not pending_reports:
        st.info("No reports pending approval")
    else:
        st.metric("Reports Pending Approval", len(pending_reports))
        
        for report_summary in pending_reports:
            with st.expander(f"{report_summary['report_name']} - {report_summary['report_id']}"):
                col1, col2 = st.columns(2)
                with col1:
                    st.write(f"**Type**: {report_summary['report_type']}")
                    st.write(f"**Generated By**: {report_summary['generated_by']}")
                    # Format dates
                    gen_at = report_summary.get('generated_at', '')
                    data_as_of = report_summary.get('data_as_of', '')
                    if gen_at and hasattr(gen_at, 'strftime'):
                        gen_at = gen_at.strftime("%m/%d/%Y")
                    if data_as_of and hasattr(data_as_of, 'strftime'):
                        data_as_of = data_as_of.strftime("%m/%d/%Y")
                    st.write(f"**Generated At**: {gen_at}")
                with col2:
                    st.write(f"**Data As Of**: {data_as_of}")
                    st.write(f"**Version**: {report_summary['version']}")
                
                # Approval actions
                col1, col2 = st.columns(2)
                with col1:
                    if st.button(f"✅ Approve", key=f"approve_{report_summary['report_id']}"):
                        approved_by = st.text_input("Approved By", value="Current User", key=f"approver_{report_summary['report_id']}")
                        if reporting.approve_report(report_summary['report_id'], approved_by):
                            st.success("Report approved!")
                            st.rerun()
                with col2:
                    if st.button(f"❌ Reject", key=f"reject_{report_summary['report_id']}"):
                        st.info("Rejection workflow would be implemented here")

with tab4:
    st.header("Submissions")
    
    # Get approved reports ready for submission
    approved_reports = reporting.list_reports(status=ReportStatus.APPROVED)
    submitted_reports = reporting.list_reports(status=ReportStatus.SUBMITTED)
    
    st.subheader("Ready for Submission")
    if not approved_reports:
        st.info("No reports ready for submission")
    else:
        for report_summary in approved_reports:
            with st.expander(f"{report_summary['report_name']} - {report_summary['report_id']}"):
                st.write(f"**Type**: {report_summary['report_type']}")
                st.write(f"**Approved By**: {report_summary.get('approved_by', 'N/A')}")
                st.write(f"**Approved At**: {report_summary.get('approved_at', 'N/A')}")
                
                if st.button(f"📤 Submit Report", key=f"submit_{report_summary['report_id']}"):
                    submitted_by = st.text_input("Submitted By", value="Current User", key=f"submitter_{report_summary['report_id']}")
                    if reporting.submit_report(report_summary['report_id'], submitted_by):
                        st.success("Report submitted!")
                        st.rerun()
    
    st.subheader("Submitted Reports")
    if not submitted_reports:
        st.info("No submitted reports")
    else:
        submitted_df = pd.DataFrame(submitted_reports)
        st.dataframe(submitted_df, use_container_width=True, hide_index=True)

with tab5:
    st.header("Audit Trail")
    
    # Select report
    all_reports = reporting.list_reports()
    if not all_reports:
        st.info("No reports available")
    else:
        selected_report_id = st.selectbox(
            "Select Report",
            [r['report_id'] for r in all_reports]
        )
        
        if selected_report_id:
            audit_trail = reporting.get_report_audit_trail(selected_report_id)
            
            if audit_trail:
                st.metric("Audit Events", len(audit_trail))
                
                # Display audit trail
                for event in audit_trail:
                    with st.expander(f"{event['event_type']} - {event['timestamp']}"):
                        st.write(f"**User**: {event['user']}")
                        st.write(f"**Timestamp**: {event['timestamp']}")
                        st.write(f"**Event Type**: {event['event_type']}")
                        if event.get('details'):
                            st.write("**Details**:")
                            st.json(event['details'])
            else:
                st.info("No audit trail events found")

# Sidebar value proposition - at bottom (if sidebar exists)
try:
    from utils.value_proposition import render_sidebar_value_proposition
    render_sidebar_value_proposition()
except:
    pass

# Footer sections - desktop full text, mobile abbreviated
render_page_footer()  # Main content footer

